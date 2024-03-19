import openpyxl
from pptx import Presentation
from pptx.util import Inches

def convert_excel_to_powerpoint(excel_file, ppt_file):
    # Load the Excel workbook
    wb = openpyxl.load_workbook(excel_file)
    sheet = wb.active
    
    # Create a new PowerPoint presentation
    prs = Presentation()
    
    # Add a slide to the presentation
    slide_layout = prs.slide_layouts[1]  # Choose a slide layout (e.g., Title and Content)
    slide = prs.slides.add_slide(slide_layout)
    
    # Get the title and content placeholders in the slide
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    # Set the title of the slide
    title_shape.text = "Excel Data to PowerPoint"
    
    # Extract data from Excel and add it to the content placeholder
    content = ""
    for row in sheet.iter_rows(values_only=True):
        content += "\n".join(str(cell) for cell in row) + "\n\n"

    content_shape.text = content
    
    # Save the PowerPoint presentation
    prs.save(ppt_file)
     
    print(f"Excel data has been converted and saved to '{ppt_file}'")

if __name__ == "__main__":
    excel_file_path = r"C:\Users\Admin\Desktop\expo\Test.xlsx"
    ppt_file_path = r"C:\Users\Admin\Desktop\expo\13march3.pptx"
    convert_excel_to_powerpoint(excel_file_path, ppt_file_path)
